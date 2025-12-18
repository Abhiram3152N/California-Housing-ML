from pptx import Presentation
from pptx.util import Inches, Pt
import os

print("Creating PowerPoint Presentation...")

# Create the presentation object
prs = Presentation()

def add_slide(prs, title_text, content_text_list):
    """Helper function to add a standard bullet point slide."""
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for point in content_text_list:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.space_after = Pt(10)
    return slide

# --- SLIDE 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide layout
slide.shapes.title.text = "Optimizing Regression Models"
slide.placeholders[1].text = "Analysis of California Housing Data\nGenerated via Python"

# --- SLIDE 2: Objective ---
add_slide(prs, "Project Objective", [
    "Goal: Predict housing prices using complex features.",
    "Method: Feature Engineering + Regularization.",
    "Models: Linear, Ridge, Lasso, ElasticNet."
])

# --- SLIDE 3: Feature Engineering ---
add_slide(prs, "Feature Engineering", [
    "1. Categorical: One-Hot Encoding on Location Zones.",
    "2. Numerical: Polynomial Features (Degree 2) & Interactions.",
    "Result: Created complex non-linear relationships for the model to learn."
])

# --- SLIDE 4: Methodology ---
add_slide(prs, "Models Tested", [
    "1. Linear Regression (Baseline): Unconstrained.",
    "2. Ridge (L2): Shrinks coefficients to handle multicollinearity.",
    "3. Lasso (L1): Performs feature selection.",
    "4. ElasticNet: Combination of L1 and L2 penalties."
])

# --- SLIDE 5: Results (Text + Table Image) ---
slide = add_slide(prs, "Model Performance Results", [
    "Linear Regression had the best Test R2 Score.",
    "Regularization (Ridge) was not strictly needed here.",
    "See table below:"
])

# Insert Table Image
img_path = 'results_table.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), height=Inches(2.5))
else:
    print(f"Warning: {img_path} not found. Slide 5 will be missing the table.")

# --- SLIDE 6: Visual Analysis (Text + Plot Image) ---
slide = add_slide(prs, "Coefficient Shrinking", [
    "Visualizing how models handle feature weights:",
    "- Linear: High variance weights.",
    "- Lasso: Drives weights to zero (Feature Selection)."
])

# Insert Plot Image
img_path = 'coefficients_plot.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), height=Inches(3.2))
else:
    print(f"Warning: {img_path} not found. Slide 6 will be missing the plot.")

# --- SLIDE 7: Conclusion ---
add_slide(prs, "Final Conclusion", [
    "Best Model: Standard Linear Regression.",
    "Reason: High generalization score (R2 ~0.61) without overfitting.",
    "Recommendation: Use Linear Regression for this specific feature set."
])

# Save the file
output_file = 'Housing_Project_Presentation.pptx'
prs.save(output_file)
print(f"SUCCESS! Presentation saved as: {output_file}")